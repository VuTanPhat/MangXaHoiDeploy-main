from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Create a presentation object
prs = Presentation()

# Set slide dimensions
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

# Slide 1: Cover Page
slide1_layout = prs.slide_layouts[6]  # Blank layout
slide1 = prs.slides.add_slide(slide1_layout)

# Add background color to cover page
background = slide1.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = RGBColor(41, 128, 185)  # Blue background

# Add title to cover page
title_box = slide1.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(2))
title_frame = title_box.text_frame
title_frame.text = "Hello World"
title_paragraph = title_frame.paragraphs[0]
title_paragraph.font.size = Pt(66)
title_paragraph.font.bold = True
title_paragraph.font.color.rgb = RGBColor(255, 255, 255)  # White text
title_paragraph.alignment = PP_ALIGN.CENTER

# Add subtitle
subtitle_box = slide1.shapes.add_textbox(Inches(1), Inches(4.5), Inches(8), Inches(1))
subtitle_frame = subtitle_box.text_frame
subtitle_frame.text = "A Simple Presentation"
subtitle_paragraph = subtitle_frame.paragraphs[0]
subtitle_paragraph.font.size = Pt(28)
subtitle_paragraph.font.color.rgb = RGBColor(255, 255, 255)
subtitle_paragraph.alignment = PP_ALIGN.CENTER

# Slide 2: Content Page
slide2_layout = prs.slide_layouts[6]  # Blank layout
slide2 = prs.slides.add_slide(slide2_layout)

# Add background color to content page
background2 = slide2.background
fill2 = background2.fill
fill2.solid()
fill2.fore_color.rgb = RGBColor(236, 240, 241)  # Light gray background

# Add title
content_title = slide2.shapes.add_textbox(Inches(1), Inches(0.5), Inches(8), Inches(1))
content_title_frame = content_title.text_frame
content_title_frame.text = "What is Hello World?"
content_title_paragraph = content_title_frame.paragraphs[0]
content_title_paragraph.font.size = Pt(44)
content_title_paragraph.font.bold = True
content_title_paragraph.font.color.rgb = RGBColor(41, 128, 185)
content_title_paragraph.alignment = PP_ALIGN.CENTER

# Add content
content_box = slide2.shapes.add_textbox(Inches(1.5), Inches(2), Inches(7), Inches(4))
content_frame = content_box.text_frame
content_frame.word_wrap = True

points = [
    "Hello World is the first program most developers learn",
    "It's a simple program that outputs 'Hello World'",
    "Used as an introduction to programming languages",
    "Symbolizes the start of a programming journey"
]

for i, point in enumerate(points):
    if i == 0:
        p = content_frame.paragraphs[0]
    else:
        p = content_frame.add_paragraph()
    p.text = "• " + point
    p.font.size = Pt(24)
    p.font.color.rgb = RGBColor(52, 73, 94)
    p.space_before = Pt(12)

# Slide 3: Closing Page
slide3_layout = prs.slide_layouts[6]  # Blank layout
slide3 = prs.slides.add_slide(slide3_layout)

# Add background color to closing page
background3 = slide3.background
fill3 = background3.fill
fill3.solid()
fill3.fore_color.rgb = RGBColor(52, 152, 219)  # Darker blue

# Add closing message
closing_box = slide3.shapes.add_textbox(Inches(1), Inches(3), Inches(8), Inches(2))
closing_frame = closing_box.text_frame
closing_frame.word_wrap = True
closing_frame.text = "Thank You!"
closing_paragraph = closing_frame.paragraphs[0]
closing_paragraph.font.size = Pt(60)
closing_paragraph.font.bold = True
closing_paragraph.font.color.rgb = RGBColor(255, 255, 255)
closing_paragraph.alignment = PP_ALIGN.CENTER

# Add closing subtitle
closing_subtitle = slide3.shapes.add_textbox(Inches(1), Inches(5), Inches(8), Inches(1))
closing_subtitle_frame = closing_subtitle.text_frame
closing_subtitle_frame.text = "Now go build amazing things!"
closing_subtitle_paragraph = closing_subtitle_frame.paragraphs[0]
closing_subtitle_paragraph.font.size = Pt(24)
closing_subtitle_paragraph.font.color.rgb = RGBColor(255, 255, 255)
closing_subtitle_paragraph.alignment = PP_ALIGN.CENTER

# Save the presentation
prs.save('Hello_World_Presentation.pptx')
print("PowerPoint presentation created successfully: Hello_World_Presentation.pptx")
